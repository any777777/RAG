#!/usr/bin/env python3
"""Build a structured Output package from the Slides folder.

The raw source folder is treated as read-only. All generated material is written
under <source-folder>/Output.
"""
from __future__ import annotations

import argparse
import hashlib
import html
import json
import math
import mimetypes
import re
import shutil
import sys
import textwrap
import time
import zipfile
from collections import defaultdict
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import fitz
from PIL import Image, ImageDraw, ImageOps

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - dependency exists in this workspace
    Presentation = None


OUTPUT_DIRNAME = "Output"
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}
PDF_EXTS = {".pdf"}
PPTX_EXTS = {".pptx"}
PPT_EXTS = {".ppt"}
TEXT_EXTS = {".md", ".txt"}
SUPPORTED_EXTS = PDF_EXTS | PPTX_EXTS | PPT_EXTS | IMAGE_EXTS | TEXT_EXTS


@dataclass
class SourceFile:
    path: Path
    rel: str
    ext: str
    size: int
    sha256: str
    source_type: str
    content_id: str = ""
    topics: list[str] = field(default_factory=list)


@dataclass
class AssetRecord:
    asset_id: str
    source_id: str
    path: str
    kind: str
    page_or_slide: str
    bytes: int
    notes: str = ""


@dataclass
class ReviewItem:
    source_id: str
    source_file: str
    reason: str
    priority: str = "medium"


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def slug(text: str, fallback: str = "source") -> str:
    value = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return value[:72] or fallback


def now_utc() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def yaml_scalar(value: Any) -> str:
    if value is None:
        return "null"
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, (int, float)):
        return str(value)
    text = str(value).replace("\\", "/").replace('"', '\\"')
    return f'"{text}"'


def frontmatter(fields: dict[str, Any]) -> str:
    lines = ["---"]
    for key, value in fields.items():
        if isinstance(value, list):
            lines.append(f"{key}:")
            if value:
                for item in value:
                    lines.append(f"  - {yaml_scalar(item)}")
            else:
                lines.append("  []")
        else:
            lines.append(f"{key}: {yaml_scalar(value)}")
    lines.append("---")
    return "\n".join(lines) + "\n\n"


def md_table(headers: list[str], rows: list[list[Any]]) -> str:
    def cell(value: Any) -> str:
        text = "" if value is None else str(value)
        text = text.replace("\r", " ").replace("\n", " ").replace("|", "\\|")
        return text

    out = ["| " + " | ".join(headers) + " |"]
    out.append("| " + " | ".join(["---"] * len(headers)) + " |")
    for row in rows:
        out.append("| " + " | ".join(cell(v) for v in row) + " |")
    return "\n".join(out)


def source_type(ext: str) -> str:
    if ext == ".pdf":
        return "pdf"
    if ext == ".pptx":
        return "pptx"
    if ext == ".ppt":
        return "ppt"
    if ext in IMAGE_EXTS:
        return "image"
    if ext in TEXT_EXTS:
        return "text"
    return "other"


def classify_topics(text: str) -> list[str]:
    t = text.lower()
    topics: list[str] = []
    rules = [
        ("Satellite Communications", [
            "satellite", "cme576", "orbital", "orbit", "look angle", "dbs",
            "uplink", "downlink", "earth station", "pratt", "sattilite",
            "سات", "ساتلايت", "satellite communication",
        ]),
        ("Machine Learning", [
            "_ml_", "machine learning", "neural", "backpropagation", "rbf",
            "som", "fuzzy", "mlp", "nn", "neural networks",
        ]),
        ("Exams and Questions", [
            "exam", "key", "questions", "tutorial", "formula sheet",
            "mid", "second", "first", "سنوات", "فيرست", "سكند",
        ]),
        ("References and Textbooks", [
            "book", "reference", "references", "ebook", "textbook",
            "كتاب", "المراجع", "t_pratt", "tlFebook".lower(),
        ]),
        ("Lecturer Notes and Photos", [
            "whiteboard", "whatsapp image", "lecturer", "photos",
            "شرح", "صور",
        ]),
    ]
    for topic, needles in rules:
        if any(n in t for n in needles):
            topics.append(topic)
    return topics or ["Unclassified"]


def ensure_output(root: Path) -> dict[str, Path]:
    out = root / OUTPUT_DIRNAME
    dirs = {
        "out": out,
        "sources": out / "sources",
        "topics": out / "topics",
        "assets": out / "assets",
        "unresolved": out / "unresolved",
        "docling_json": out / "docling-json",
        "rag_chunks": out / "rag-chunks",
    }
    for path in dirs.values():
        path.mkdir(parents=True, exist_ok=True)
    return dirs


def iter_raw_files(root: Path) -> list[Path]:
    out = root / OUTPUT_DIRNAME
    files: list[Path] = []
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        if out in path.parents:
            continue
        files.append(path)
    return files


def inventory(root: Path) -> list[SourceFile]:
    records: list[SourceFile] = []
    for path in iter_raw_files(root):
        ext = path.suffix.lower()
        rel = path.relative_to(root).as_posix()
        digest = sha256(path)
        records.append(
            SourceFile(
                path=path,
                rel=rel,
                ext=ext,
                size=path.stat().st_size,
                sha256=digest,
                source_type=source_type(ext),
                topics=classify_topics(rel),
            )
        )
    used_ids: set[str] = set()
    for rec in records:
        base = slug(rec.path.stem, fallback=f"source-{rec.sha256[:10]}")
        cid = f"{base}-{rec.sha256[:10]}"
        if cid in used_ids:
            cid = f"{cid}-{len(used_ids)}"
        used_ids.add(cid)
        rec.content_id = cid
    return records


def canonical_by_hash(records: list[SourceFile]) -> dict[str, SourceFile]:
    groups: dict[str, list[SourceFile]] = defaultdict(list)
    for rec in records:
        groups[rec.sha256].append(rec)
    canonical: dict[str, SourceFile] = {}
    for digest, group in groups.items():
        canonical[digest] = sorted(group, key=lambda r: (r.rel.count("/"), r.rel.lower()))[0]
    return canonical


def write_json(path: Path, data: dict[str, Any]) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def relative_to_output(path: Path, out: Path) -> str:
    return path.relative_to(out).as_posix()


def make_contact_sheet(asset_paths: list[Path], out_path: Path, title: str, max_items: int = 40) -> None:
    images: list[tuple[Path, Image.Image]] = []
    for path in asset_paths[:max_items]:
        try:
            img = Image.open(path).convert("RGB")
            img.thumbnail((220, 160))
            thumb = ImageOps.expand(img, border=1, fill=(200, 200, 200))
            images.append((path, thumb))
        except Exception:
            continue
    if not images:
        return
    cols = 4
    rows = math.ceil(len(images) / cols)
    cell_w, cell_h = 260, 210
    header_h = 44
    sheet = Image.new("RGB", (cols * cell_w, rows * cell_h + header_h), "white")
    draw = ImageDraw.Draw(sheet)
    draw.text((12, 12), title[:160], fill=(0, 0, 0))
    for idx, (path, img) in enumerate(images):
        row, col = divmod(idx, cols)
        x = col * cell_w + 10
        y = header_h + row * cell_h + 10
        sheet.paste(img, (x, y))
        draw.text((x, y + 166), path.name[:36], fill=(0, 0, 0))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path)


def extract_pdf(
    rec: SourceFile,
    dirs: dict[str, Path],
    max_captures_per_source: int,
) -> tuple[str, dict[str, Any], list[AssetRecord], list[ReviewItem]]:
    assets: list[AssetRecord] = []
    reviews: list[ReviewItem] = []
    out = dirs["out"]
    asset_dir = dirs["assets"] / rec.content_id / "extracted"
    asset_dir.mkdir(parents=True, exist_ok=True)
    markdown_parts: list[str] = []
    pages: list[dict[str, Any]] = []
    extracted_asset_paths: list[Path] = []
    capture_count = 0

    try:
        doc = fitz.open(rec.path)
    except Exception as exc:
        reviews.append(ReviewItem(rec.content_id, rec.rel, f"PDF open failed: {exc}", "high"))
        return "", {"error": str(exc), "pages": []}, assets, reviews

    metadata = dict(doc.metadata or {})
    for idx, page in enumerate(doc, 1):
        try:
            text = page.get_text("text").strip()
        except Exception as exc:
            text = ""
            reviews.append(ReviewItem(rec.content_id, rec.rel, f"page {idx}: text extraction failed: {exc}", "medium"))
        is_sparse = len(text) < 80
        page_record: dict[str, Any] = {
            "page": idx,
            "text_characters": len(text),
            "sparse_text": is_sparse,
            "assets": [],
        }
        markdown_parts.append(f"## Page {idx}\n\n{text or '_No reliable embedded text extracted._'}")
        if is_sparse:
            reviews.append(ReviewItem(rec.content_id, rec.rel, f"page {idx}: sparse embedded text; OCR/visual review needed", "medium"))
        try:
            page_images = page.get_images(full=True)
        except Exception:
            page_images = []
        for j, img in enumerate(page_images, 1):
            try:
                raw = doc.extract_image(img[0])
                ext = raw.get("ext", "png")
                asset_path = asset_dir / f"page-{idx:04d}-image-{j:03d}.{ext}"
                asset_path.write_bytes(raw["image"])
                extracted_asset_paths.append(asset_path)
                asset = AssetRecord(
                    asset_id=f"{rec.content_id}-p{idx:04d}-img{j:03d}",
                    source_id=rec.content_id,
                    path=relative_to_output(asset_path, out),
                    kind="unknown",
                    page_or_slide=f"page {idx}",
                    bytes=asset_path.stat().st_size,
                    notes="embedded PDF image",
                )
                assets.append(asset)
                page_record["assets"].append(asset.asset_id)
            except Exception as exc:
                reviews.append(ReviewItem(rec.content_id, rec.rel, f"page {idx}: image extraction failed: {exc}", "low"))
        if (is_sparse or idx == 1) and capture_count < max_captures_per_source:
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0), alpha=False)
                capture_path = asset_dir / f"page-{idx:04d}-capture.png"
                pix.save(capture_path)
                capture_count += 1
                extracted_asset_paths.append(capture_path)
                asset = AssetRecord(
                    asset_id=f"{rec.content_id}-p{idx:04d}-capture",
                    source_id=rec.content_id,
                    path=relative_to_output(capture_path, out),
                    kind="page-capture",
                    page_or_slide=f"page {idx}",
                    bytes=capture_path.stat().st_size,
                    notes="rendered page capture for layout/OCR review",
                )
                assets.append(asset)
                page_record["assets"].append(asset.asset_id)
            except Exception as exc:
                reviews.append(ReviewItem(rec.content_id, rec.rel, f"page {idx}: page capture failed: {exc}", "low"))
        pages.append(page_record)

    if capture_count >= max_captures_per_source:
        reviews.append(
            ReviewItem(
                rec.content_id,
                rec.rel,
                f"page capture limit reached ({max_captures_per_source}); inspect remaining sparse pages before OCR",
                "medium",
            )
        )
    contact = dirs["assets"] / rec.content_id / "contact-sheet.png"
    make_contact_sheet(extracted_asset_paths, contact, rec.rel)
    if contact.exists():
        assets.append(
            AssetRecord(
                asset_id=f"{rec.content_id}-contact-sheet",
                source_id=rec.content_id,
                path=relative_to_output(contact, out),
                kind="page-capture",
                page_or_slide="multi",
                bytes=contact.stat().st_size,
                notes="contact sheet for visual review",
            )
        )
    structured = {
        "source_id": rec.content_id,
        "source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "metadata": metadata,
        "page_count": len(pages),
        "pages": pages,
        "extraction_method": "pymupdf",
    }
    return "\n\n".join(markdown_parts), structured, assets, reviews


def iter_shape_text(shape: Any) -> list[str]:
    texts: list[str] = []
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        value = (shape.text or "").strip()
        if value:
            texts.append(value)
    if hasattr(shape, "has_table") and shape.has_table:
        rows = []
        for row in shape.table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            texts.append("\n".join(rows))
    return texts


def extract_pptx(rec: SourceFile, dirs: dict[str, Path]) -> tuple[str, dict[str, Any], list[AssetRecord], list[ReviewItem]]:
    assets: list[AssetRecord] = []
    reviews: list[ReviewItem] = []
    out = dirs["out"]
    asset_dir = dirs["assets"] / rec.content_id / "extracted"
    asset_dir.mkdir(parents=True, exist_ok=True)
    markdown_parts: list[str] = []
    slides: list[dict[str, Any]] = []
    extracted_asset_paths: list[Path] = []

    if Presentation is None:
        reviews.append(ReviewItem(rec.content_id, rec.rel, "python-pptx unavailable; PPTX not extracted", "high"))
        return "", {"error": "python-pptx unavailable"}, assets, reviews

    try:
        prs = Presentation(str(rec.path))
    except Exception as exc:
        reviews.append(ReviewItem(rec.content_id, rec.rel, f"PPTX open failed: {exc}", "high"))
        return "", {"error": str(exc)}, assets, reviews

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        slide_asset_ids: list[str] = []
        for shape_idx, shape in enumerate(slide.shapes, 1):
            slide_texts.extend(iter_shape_text(shape))
            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    asset_path = asset_dir / f"slide-{slide_idx:04d}-image-{shape_idx:03d}.{ext}"
                    asset_path.write_bytes(image.blob)
                    extracted_asset_paths.append(asset_path)
                    asset = AssetRecord(
                        asset_id=f"{rec.content_id}-s{slide_idx:04d}-img{shape_idx:03d}",
                        source_id=rec.content_id,
                        path=relative_to_output(asset_path, out),
                        kind="unknown",
                        page_or_slide=f"slide {slide_idx}",
                        bytes=asset_path.stat().st_size,
                        notes="PPTX image shape",
                    )
                    assets.append(asset)
                    slide_asset_ids.append(asset.asset_id)
                except Exception as exc:
                    reviews.append(ReviewItem(rec.content_id, rec.rel, f"slide {slide_idx}: image extraction failed: {exc}", "low"))
        text = "\n\n".join(t for t in slide_texts if t).strip()
        if not text and not slide_asset_ids:
            reviews.append(ReviewItem(rec.content_id, rec.rel, f"slide {slide_idx}: no text/assets extracted; visual review needed", "medium"))
        markdown_parts.append(f"## Slide {slide_idx}\n\n{text or '_No text extracted._'}")
        slides.append({
            "slide": slide_idx,
            "text_characters": len(text),
            "assets": slide_asset_ids,
        })

    # Fallback ZIP media extraction catches media not exposed as picture shapes.
    try:
        with zipfile.ZipFile(rec.path) as archive:
            for item in archive.namelist():
                if not item.startswith("ppt/media/") or item.endswith("/"):
                    continue
                asset_path = asset_dir / Path(item).name
                if not asset_path.exists():
                    asset_path.write_bytes(archive.read(item))
                    extracted_asset_paths.append(asset_path)
                    assets.append(
                        AssetRecord(
                            asset_id=f"{rec.content_id}-media-{slug(Path(item).stem, Path(item).stem)}",
                            source_id=rec.content_id,
                            path=relative_to_output(asset_path, out),
                            kind="unknown",
                            page_or_slide="unknown slide",
                            bytes=asset_path.stat().st_size,
                            notes="PPTX package media fallback",
                        )
                    )
    except Exception as exc:
        reviews.append(ReviewItem(rec.content_id, rec.rel, f"PPTX package media scan failed: {exc}", "low"))

    contact = dirs["assets"] / rec.content_id / "contact-sheet.png"
    make_contact_sheet(extracted_asset_paths, contact, rec.rel)
    if contact.exists():
        assets.append(
            AssetRecord(
                asset_id=f"{rec.content_id}-contact-sheet",
                source_id=rec.content_id,
                path=relative_to_output(contact, out),
                kind="page-capture",
                page_or_slide="multi",
                bytes=contact.stat().st_size,
                notes="contact sheet for visual review",
            )
        )

    structured = {
        "source_id": rec.content_id,
        "source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "slide_count": len(slides),
        "slides": slides,
        "extraction_method": "python-pptx",
    }
    return "\n\n".join(markdown_parts), structured, assets, reviews


def extract_image(rec: SourceFile, dirs: dict[str, Path]) -> tuple[str, dict[str, Any], list[AssetRecord], list[ReviewItem]]:
    out = dirs["out"]
    asset_dir = dirs["assets"] / rec.content_id / "extracted"
    asset_dir.mkdir(parents=True, exist_ok=True)
    dest = asset_dir / rec.path.name
    shutil.copy2(rec.path, dest)
    width = height = None
    try:
        with Image.open(dest) as img:
            width, height = img.size
    except Exception:
        pass
    contact = dirs["assets"] / rec.content_id / "contact-sheet.png"
    make_contact_sheet([dest], contact, rec.rel)
    assets = [
        AssetRecord(
            asset_id=f"{rec.content_id}-image",
            source_id=rec.content_id,
            path=relative_to_output(dest, out),
            kind="photo",
            page_or_slide="image",
            bytes=dest.stat().st_size,
            notes="raw image evidence copied",
        )
    ]
    if contact.exists():
        assets.append(
            AssetRecord(
                asset_id=f"{rec.content_id}-contact-sheet",
                source_id=rec.content_id,
                path=relative_to_output(contact, out),
                kind="page-capture",
                page_or_slide="image",
                bytes=contact.stat().st_size,
                notes="contact sheet for visual review",
            )
        )
    markdown = "_Image evidence copied. OCR was not run in this local first pass._"
    structured = {
        "source_id": rec.content_id,
        "source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "width": width,
        "height": height,
        "mime_type": mimetypes.guess_type(rec.path.name)[0],
        "extraction_method": "copy-image",
    }
    reviews = [ReviewItem(rec.content_id, rec.rel, "image OCR and visual classification required", "medium")]
    return markdown, structured, assets, reviews


def extract_text_file(rec: SourceFile) -> tuple[str, dict[str, Any], list[AssetRecord], list[ReviewItem]]:
    try:
        text = rec.path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        text = ""
        reviews = [ReviewItem(rec.content_id, rec.rel, f"text read failed: {exc}", "medium")]
    else:
        reviews = []
    structured = {
        "source_id": rec.content_id,
        "source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "text_characters": len(text),
        "extraction_method": "read-text",
    }
    return text, structured, [], reviews


def unsupported_ppt(rec: SourceFile) -> tuple[str, dict[str, Any], list[AssetRecord], list[ReviewItem]]:
    markdown = (
        "_Legacy `.ppt` file cataloged but not extracted. LibreOffice/PowerPoint conversion "
        "to `.pptx` or PDF is required before structured extraction._"
    )
    structured = {
        "source_id": rec.content_id,
        "source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "extraction_method": "catalog-only",
        "conversion_needed": True,
    }
    reviews = [ReviewItem(rec.content_id, rec.rel, "legacy .ppt requires conversion to PPTX/PDF; LibreOffice not available", "high")]
    return markdown, structured, [], reviews


def split_chunks(text: str, target_chars: int = 2500) -> list[str]:
    clean = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not clean:
        return []
    parts: list[str] = []
    current: list[str] = []
    size = 0
    for paragraph in clean.split("\n\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        if current and size + len(paragraph) > target_chars:
            parts.append("\n\n".join(current))
            current, size = [], 0
        current.append(paragraph)
        size += len(paragraph)
    if current:
        parts.append("\n\n".join(current))
    return parts


def write_source_record(
    rec: SourceFile,
    aliases: list[SourceFile],
    markdown: str,
    structured: dict[str, Any],
    dirs: dict[str, Path],
    topics: list[str],
    assets: list[AssetRecord],
    reviews: list[ReviewItem],
) -> None:
    fields = {
        "source_id": rec.content_id,
        "canonical_source_file": rec.rel,
        "source_type": rec.source_type,
        "sha256": rec.sha256,
        "size_bytes": rec.size,
        "status": "needs-review" if reviews else "verified-first-pass",
        "topics": topics,
        "duplicate_source_files": [a.rel for a in aliases if a.rel != rec.rel],
        "extraction_method": structured.get("extraction_method"),
        "generated_at": now_utc(),
    }
    body = [
        f"# {rec.path.stem}",
        "",
        "## Source Files",
        "",
        md_table(["Role", "Path", "Size bytes"], [["canonical", rec.rel, rec.size]] + [["duplicate", a.rel, a.size] for a in aliases if a.rel != rec.rel]),
        "",
        "## Topics",
        "",
        "\n".join(f"- {topic}" for topic in topics),
        "",
        "## Extracted Asset Count",
        "",
        str(len(assets)),
        "",
        "## Review Notes",
        "",
        "\n".join(f"- [{item.priority}] {item.reason}" for item in reviews) or "_No review notes in first pass._",
        "",
        "## Extracted Text",
        "",
        markdown or "_No extracted text._",
    ]
    path = dirs["sources"] / f"{rec.content_id}.md"
    path.write_text(frontmatter(fields) + "\n".join(body).rstrip() + "\n", encoding="utf-8")


def write_rag_chunks(rec: SourceFile, markdown: str, dirs: dict[str, Path], topics: list[str]) -> int:
    chunks = split_chunks(markdown)
    if not chunks:
        return 0
    chunk_dir = dirs["rag_chunks"] / rec.content_id
    chunk_dir.mkdir(parents=True, exist_ok=True)
    for idx, chunk in enumerate(chunks, 1):
        fields = {
            "chunk_id": f"{rec.content_id}-chunk-{idx:04d}",
            "source_id": rec.content_id,
            "source_file": rec.rel,
            "source_type": rec.source_type,
            "topics": topics,
            "chunk_index": idx,
            "confidence": "first-pass",
            "extraction_method": "structured-local",
        }
        (chunk_dir / f"chunk-{idx:04d}.md").write_text(frontmatter(fields) + chunk + "\n", encoding="utf-8")
    return len(chunks)


def maybe_run_docling(
    rec: SourceFile,
    dirs: dict[str, Path],
    mode: str,
    max_mb: float,
    budget: dict[str, int],
    reviews: list[ReviewItem],
) -> None:
    if mode == "off" or rec.ext not in {".pdf", ".pptx"}:
        return
    if rec.size > max_mb * 1024 * 1024:
        reviews.append(ReviewItem(rec.content_id, rec.rel, f"Docling skipped: file larger than {max_mb:g} MB threshold", "low"))
        return
    if budget["remaining"] <= 0:
        reviews.append(ReviewItem(rec.content_id, rec.rel, "Docling skipped: per-run Docling file budget reached", "low"))
        return
    budget["remaining"] -= 1
    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(str(rec.path))
        doc = result.document
        docling_data = doc.export_to_dict()
        docling_md = doc.export_to_markdown()
        (dirs["docling_json"] / f"{rec.content_id}.docling.md").write_text(docling_md, encoding="utf-8")
        write_json(dirs["docling_json"] / f"{rec.content_id}.docling.json", docling_data)
    except Exception as exc:
        reviews.append(ReviewItem(rec.content_id, rec.rel, f"Docling conversion failed: {exc}", "medium"))


def write_indexes(
    root: Path,
    dirs: dict[str, Path],
    records: list[SourceFile],
    canonical: dict[str, SourceFile],
    all_assets: list[AssetRecord],
    all_reviews: list[ReviewItem],
    chunk_counts: dict[str, int],
    started_at: str,
    finished_at: str,
) -> None:
    out = dirs["out"]
    groups: dict[str, list[SourceFile]] = defaultdict(list)
    for rec in records:
        groups[rec.sha256].append(rec)

    catalog_rows = []
    for rec in sorted(records, key=lambda r: r.rel.lower()):
        canon = canonical[rec.sha256]
        duplicate_of = "" if rec.rel == canon.rel else canon.content_id
        catalog_rows.append([
            rec.rel,
            rec.source_type,
            rec.size,
            rec.sha256[:16],
            rec.content_id,
            duplicate_of,
            ", ".join(rec.topics),
        ])
    (out / "00-catalog.md").write_text(
        frontmatter({"title": "Slides Catalog", "generated_at": finished_at, "source_folder": root.as_posix(), "file_count": len(records)})
        + "# Slides Catalog\n\n"
        + md_table(["Source file", "Type", "Size bytes", "SHA-256 prefix", "Source ID", "Duplicate of", "Topics"], catalog_rows)
        + "\n",
        encoding="utf-8",
    )

    register_rows = []
    for digest, group in sorted(groups.items(), key=lambda item: canonical[item[0]].rel.lower()):
        canon = canonical[digest]
        register_rows.append([
            canon.content_id,
            canon.rel,
            canon.source_type,
            digest,
            len(group),
            ", ".join(canon.topics),
            "needs-review",
        ])
    (out / "02-source-register.md").write_text(
        frontmatter({"title": "Source Register", "generated_at": finished_at, "unique_content_count": len(groups)})
        + "# Source Register\n\n"
        + md_table(["Source ID", "Canonical file", "Type", "SHA-256", "Copies", "Topics", "Status"], register_rows)
        + "\n",
        encoding="utf-8",
    )

    asset_rows = [
        [a.asset_id, a.source_id, a.kind, a.page_or_slide, a.bytes, a.path, a.notes]
        for a in sorted(all_assets, key=lambda item: (item.source_id, item.path))
    ]
    (out / "03-assets-index.md").write_text(
        frontmatter({"title": "Assets Index", "generated_at": finished_at, "asset_count": len(all_assets)})
        + "# Assets Index\n\n"
        + md_table(["Asset ID", "Source ID", "Kind", "Page/Slide", "Size bytes", "Path", "Notes"], asset_rows)
        + "\n\n<!-- MANUAL START -->\n_Add reviewed classifications here: diagram, chart, table, equation, code, screenshot, photo, page-capture._\n<!-- MANUAL END -->\n",
        encoding="utf-8",
    )

    review_rows = [
        [idx, item.priority, item.source_id, item.source_file, item.reason]
        for idx, item in enumerate(all_reviews, 1)
    ]
    (out / "04-review-queue.md").write_text(
        frontmatter({"title": "Review Queue", "generated_at": finished_at, "review_item_count": len(all_reviews)})
        + "# Review Queue\n\n"
        + (md_table(["#", "Priority", "Source ID", "Source file", "Reason"], review_rows) if review_rows else "_No review items._")
        + "\n\n<!-- MANUAL START -->\n_Add manual review decisions here._\n<!-- MANUAL END -->\n",
        encoding="utf-8",
    )

    dedupe_rows = []
    for digest, group in groups.items():
        if len(group) <= 1:
            continue
        canon = canonical[digest]
        dedupe_rows.append([
            digest,
            canon.content_id,
            len(group),
            "; ".join(g.rel for g in sorted(group, key=lambda r: r.rel.lower())),
        ])
    (out / "06-deduplication-report.md").write_text(
        frontmatter({"title": "Deduplication Report", "generated_at": finished_at, "duplicate_group_count": len(dedupe_rows)})
        + "# Deduplication Report\n\n"
        + (md_table(["SHA-256", "Canonical Source ID", "Copy count", "Files"], dedupe_rows) if dedupe_rows else "_No exact duplicates found._")
        + "\n",
        encoding="utf-8",
    )

    topic_map: dict[str, list[SourceFile]] = defaultdict(list)
    for digest, group in groups.items():
        canon = canonical[digest]
        for topic in canon.topics:
            topic_map[topic].append(canon)
    topic_rows = []
    for topic, items in sorted(topic_map.items()):
        topic_slug = slug(topic, "topic")
        topic_rows.append([topic, len(items), f"topics/{topic_slug}.md"])
        topic_body = [
            f"# {topic}",
            "",
            md_table(
                ["Source ID", "Canonical file", "Type", "RAG chunks"],
                [[item.content_id, item.rel, item.source_type, chunk_counts.get(item.content_id, 0)] for item in sorted(items, key=lambda r: r.rel.lower())],
            ),
            "",
            "<!-- MANUAL START -->",
            "_Add reviewed topic ranges, prerequisites, and conflicts here._",
            "<!-- MANUAL END -->",
        ]
        (dirs["topics"] / f"{topic_slug}.md").write_text(
            frontmatter({"title": topic, "generated_at": finished_at, "source_count": len(items)})
            + "\n".join(topic_body)
            + "\n",
            encoding="utf-8",
        )
    (out / "01-topic-map.md").write_text(
        frontmatter({"title": "Topic Map", "generated_at": finished_at, "topic_count": len(topic_rows)})
        + "# Topic Map\n\n"
        + md_table(["Topic", "Unique sources", "Topic file"], topic_rows)
        + "\n\n<!-- MANUAL START -->\n_Add manual grouping decisions here._\n<!-- MANUAL END -->\n",
        encoding="utf-8",
    )

    type_counts = defaultdict(int)
    for rec in records:
        type_counts[rec.source_type] += 1
    log = [
        "# Processing Log",
        "",
        f"- Started: `{started_at}`",
        f"- Finished: `{finished_at}`",
        f"- Raw files: `{len(records)}`",
        f"- Unique content items: `{len(groups)}`",
        f"- Exact duplicate groups: `{sum(1 for group in groups.values() if len(group) > 1)}`",
        f"- Extracted assets: `{len(all_assets)}`",
        f"- Review items: `{len(all_reviews)}`",
        f"- RAG chunks: `{sum(chunk_counts.values())}`",
        "",
        "## Type Counts",
        "",
        md_table(["Type", "Count"], [[k, v] for k, v in sorted(type_counts.items())]),
    ]
    (out / "05-processing-log.md").write_text(
        frontmatter({"title": "Processing Log", "generated_at": finished_at}) + "\n".join(log) + "\n",
        encoding="utf-8",
    )

    readme = textwrap.dedent(
        f"""
        # Slides Output

        This folder is a first-pass structured package generated from the raw `Slides` folder.

        - Raw source files were not moved, renamed, or deleted.
        - Exact duplicate content was processed once and referenced through duplicate records.
        - `.ppt` files were cataloged but require conversion to `.pptx` or PDF before extraction.
        - Image OCR and difficult scanned pages remain in `04-review-queue.md`.

        Suggested next step: review `04-review-queue.md`, then classify important visual assets in the manual section of `03-assets-index.md`.
        """
    ).strip()
    (out / "README.md").write_text(frontmatter({"title": "Slides Output", "generated_at": finished_at}) + readme + "\n", encoding="utf-8")


def process(root: Path, args: argparse.Namespace) -> None:
    started_at = now_utc()
    dirs = ensure_output(root)
    records = inventory(root)
    canonical = canonical_by_hash(records)
    groups: dict[str, list[SourceFile]] = defaultdict(list)
    for rec in records:
        groups[rec.sha256].append(rec)

    all_assets: list[AssetRecord] = []
    all_reviews: list[ReviewItem] = []
    chunk_counts: dict[str, int] = {}
    docling_budget = {"remaining": args.docling_max_files}

    processed = 0
    for digest, aliases in sorted(groups.items(), key=lambda item: canonical[item[0]].rel.lower()):
        rec = canonical[digest]
        rec.topics = classify_topics(rec.rel + " " + rec.path.stem)
        print(f"[{processed + 1}/{len(groups)}] {rec.source_type}: {rec.rel}", flush=True)
        if rec.ext == ".pdf":
            markdown, structured, assets, reviews = extract_pdf(rec, dirs, args.max_captures_per_source)
            maybe_run_docling(rec, dirs, args.docling_mode, args.docling_max_mb, docling_budget, reviews)
        elif rec.ext == ".pptx":
            markdown, structured, assets, reviews = extract_pptx(rec, dirs)
            maybe_run_docling(rec, dirs, args.docling_mode, args.docling_max_mb, docling_budget, reviews)
        elif rec.ext == ".ppt":
            markdown, structured, assets, reviews = unsupported_ppt(rec)
        elif rec.ext in IMAGE_EXTS:
            markdown, structured, assets, reviews = extract_image(rec, dirs)
        elif rec.ext in TEXT_EXTS:
            markdown, structured, assets, reviews = extract_text_file(rec)
        else:
            markdown = "_Unsupported source type cataloged only._"
            structured = {"source_id": rec.content_id, "source_file": rec.rel, "source_type": rec.source_type, "sha256": rec.sha256, "extraction_method": "catalog-only"}
            assets = []
            reviews = [ReviewItem(rec.content_id, rec.rel, "unsupported source type", "medium")]

        structured["duplicate_source_files"] = [a.rel for a in aliases if a.rel != rec.rel]
        structured["topics"] = rec.topics
        write_json(dirs["docling_json"] / f"{rec.content_id}.structured.json", structured)
        write_source_record(rec, aliases, markdown, structured, dirs, rec.topics, assets, reviews)
        chunk_counts[rec.content_id] = write_rag_chunks(rec, markdown, dirs, rec.topics)
        all_assets.extend(assets)
        all_reviews.extend(reviews)
        processed += 1

    finished_at = now_utc()
    write_indexes(root, dirs, records, canonical, all_assets, all_reviews, chunk_counts, started_at, finished_at)
    print(f"[OK] Wrote structured package to {dirs['out']}", flush=True)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build Slides/Output structured extraction package.")
    parser.add_argument("source_folder", nargs="?", default="Slides", help="Source folder containing raw educational material.")
    parser.add_argument("--max-captures-per-source", type=int, default=80, help="Maximum rendered PDF page captures per source.")
    parser.add_argument("--docling-mode", choices=["off", "small"], default="small", help="Run Docling on small PDF/PPTX files only, or disable it.")
    parser.add_argument("--docling-max-mb", type=float, default=0.2, help="Maximum file size for Docling conversion when docling-mode=small.")
    parser.add_argument("--docling-max-files", type=int, default=3, help="Maximum Docling conversions per run.")
    return parser.parse_args()


def main() -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    args = parse_args()
    root = Path(args.source_folder).resolve()
    if not root.exists() or not root.is_dir():
        print(f"Source folder does not exist: {root}", file=sys.stderr)
        return 2
    start = time.time()
    process(root, args)
    elapsed = time.time() - start
    print(f"[OK] elapsed_seconds={elapsed:.1f}", flush=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
